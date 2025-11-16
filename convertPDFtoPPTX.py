import os
import fitz  # PyMuPDF
from pptx import Presentation

input_folder = r"D:\02-Training-Courses\GIT-Sharing-Udemy-Document\pdf"
output_folder = r"D:\02-Training-Courses\GIT-Sharing-Udemy-Document\pptx"
os.makedirs(output_folder, exist_ok=True)

pdf_files = [f for f in os.listdir(input_folder) if f.lower().endswith(".pdf")]

for pdf_file in pdf_files:
    pdf_path = os.path.join(input_folder, pdf_file)
    pptx_name = os.path.splitext(pdf_file)[0] + ".pptx"
    pptx_path = os.path.join(output_folder, pptx_name)

    print(f"Processing: {pdf_file}")
    
    # Mở PDF
    doc = fitz.open(pdf_path)
    prs = Presentation()
    
    for page_num, page in enumerate(doc):
        # Render page thành ảnh
        pix = page.get_pixmap(dpi=150)  # tăng dpi càng nét
        img_path = f"temp_page_{page_num}.png"
        pix.save(img_path)
        
        # Tạo slide blank
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Kích thước slide
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Tỷ lệ ảnh PDF
        img_ratio = pix.width / pix.height
        slide_ratio = slide_width / slide_height
        
        # Tính width, height và canh giữa
        if img_ratio > slide_ratio:
            # ảnh rộng hơn slide → fit theo width
            width = slide_width
            height = int(slide_width / img_ratio)
            left = 0
            top = int((slide_height - height) / 2)
        else:
            # ảnh cao hơn slide → fit theo height
            height = slide_height
            width = int(slide_height * img_ratio)
            top = 0
            left = int((slide_width - width) / 2)
        
        # Chèn ảnh vào slide
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        
        # Xoá ảnh tạm
        os.remove(img_path)
    
    # Lưu PPTX
    prs.save(pptx_path)
    print(f"Saved PPTX: {pptx_name}")

print("Done!")
